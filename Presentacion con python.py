from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, layout, title, content):
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]

    title_shape.text = title
    tf = content_shape.text_frame
    tf.text = content
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.LEFT

    return slide

def create_presentation():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Slide 1
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Transcriptor de Audio a Texto"
    subtitle.text = "Una herramienta para convertir audio en texto"

    # Slide 2
    add_slide(prs, bullet_slide_layout, "Características Principales",
              "• Manejo de múltiples formatos de audio\n"
              "• Transcripción de audio a texto\n"
              "• Traducción de transcripciones\n"
              "• Reproducción de archivos de audio\n"
              "• Exportación de transcripciones")

    # Slide 3
    add_slide(prs, bullet_slide_layout, "Manejo de Archivos de Audio",
              "• Selección múltiple de archivos\n"
              "• Visualización en lista\n"
              "• Opción para borrar archivos\n"
              "• Soporta formatos: MP3, WAV, FLAC, OGG, M4A, MP4, AAC")

    # Slide 4
    add_slide(prs, bullet_slide_layout, "Transcripción y Traducción",
              "• Transcripción de audio a texto\n"
              "• Soporte para múltiples idiomas\n"
              "• Opción de traducción a otro idioma\n"
              "• Visualización de progreso")

    # Slide 5
    add_slide(prs, bullet_slide_layout, "Reproducción de Audio",
              "• Reproducción de archivos seleccionados\n"
              "• Controles: reproducir, pausar, reanudar, detener\n"
              "• Visualización del tiempo de reproducción")

    # Slide 6
    add_slide(prs, bullet_slide_layout, "Manejo de Texto",
              "• Visualización de transcripciones\n"
              "• Opción para limpiar el área de texto\n"
              "• Exportación a archivo de texto")

    # Slide 7
    add_slide(prs, bullet_slide_layout, "Tecnologías Utilizadas",
              "• Python como lenguaje principal\n"
              "• Tkinter para la interfaz gráfica\n"
              "• Bibliotecas: speech_recognition, pydub, googletrans, pygame")

    # Slide 8
    add_slide(prs, bullet_slide_layout, "¡Gracias por su atención!",
              "El Transcriptor de Audio a Texto está diseñado para facilitar\n"
              "la conversión de audio a texto con funciones adicionales útiles.")

    prs.save('Presentacion_Transcriptor.pptx')

if __name__ == '__main__':
    create_presentation()